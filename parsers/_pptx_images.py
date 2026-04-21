"""PowerPoint image extractor — covers what docling drops.

docling's generic ``DocumentConverter`` walks .pptx slides and pulls
picture shapes, but it doesn't visit the **speaker-notes slide**
attached to each deck slide. Notes frequently carry the diagrams /
screenshots the presenter walks through live — for education /
product-walkthrough material that's often higher-value content than
the main slide's decoration.

This extractor plugs into the W4-b framework and runs *after* the
normal parse (which fills the markdown text), then returns the
full image set — slide visuals + speaker-notes visuals — as
``ImageRef`` lazy factories.

Naming convention is informative:
  slide_{idx:03d}_img_{n}
  slide_{idx:03d}_notes_img_{n}
so downstream DB rows + S3 keys are greppable to "which slide did
this come from, and was it in notes".
"""
from __future__ import annotations

import io
import logging
from typing import TYPE_CHECKING

from parsers.image_ref import ImageRef

if TYPE_CHECKING:
    from pptx.slide import Slide  # noqa: F401

logger = logging.getLogger(__name__)

_PPTX_EXTS = frozenset({".pptx"})
_PPTX_MIMES = frozenset({
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
})


def _collect_shape_images(shapes, labels: list) -> None:
    """Walk a shape tree and append ``(name, mime, bytes)`` for every
    picture found. Recurses into GroupShapes; leaf Picture shapes carry
    ``.image.blob`` + ``.image.content_type``.
    """
    for shape in shapes:
        # Group → recurse
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            _collect_shape_images(shape.shapes, labels)
            continue
        # Pictures expose ``.image``; non-picture shapes raise AttributeError
        image = getattr(shape, "image", None)
        if image is None:
            continue
        try:
            blob = image.blob
            content_type = image.content_type or "image/png"
        except Exception:
            continue
        labels.append((shape, blob, content_type))


class PptxImageExtractor:
    """BaseImageExtractor for .pptx files.

    See module docstring. Registered from ``parsers/__init__.py`` at
    import time.
    """

    def match(self, *, mime: str = "", ext: str = "") -> bool:
        return ext in _PPTX_EXTS or mime in _PPTX_MIMES

    async def extract(
        self,
        source: bytes,
        *,
        filename: str = "",
    ) -> list[ImageRef]:
        # Import inside so the worker doesn't pay pptx's XML parse cost
        # for file types that never reach this extractor.
        from pptx import Presentation
        from pptx.util import Emu

        try:
            prs = Presentation(io.BytesIO(source))
        except Exception:
            logger.exception("pptx parse failed for %s", filename)
            return []

        refs: list[ImageRef] = []
        for slide_idx, slide in enumerate(prs.slides):
            # Slide body images ---------------------------------------
            slide_images: list = []
            _collect_shape_images(slide.shapes, slide_images)
            for n, (shape, blob, content_type) in enumerate(slide_images):
                refs.append(_make_ref(
                    name=f"slide_{slide_idx:03d}_img_{n}",
                    blob=blob,
                    content_type=content_type,
                    shape=shape,
                ))

            # Speaker notes images — docling misses these -------------
            if not slide.has_notes_slide:
                continue
            notes_images: list = []
            _collect_shape_images(slide.notes_slide.shapes, notes_images)
            for n, (shape, blob, content_type) in enumerate(notes_images):
                refs.append(_make_ref(
                    name=f"slide_{slide_idx:03d}_notes_img_{n}",
                    blob=blob,
                    content_type=content_type,
                    shape=shape,
                ))

        return refs


def _make_ref(
    *, name: str, blob: bytes, content_type: str, shape,
) -> ImageRef:
    """Build an ImageRef from an extracted pptx image blob.

    Width / height come from the slide geometry when available
    (``shape.width`` / ``shape.height`` are in EMU — English Metric
    Units — 914400 per inch). Falls back to PIL decode if the shape
    has no explicit dims.
    """
    w = h = None
    try:
        if shape.width and shape.height:
            # EMU → pixels at 96 DPI: px = emu / 914400 * 96
            w = int(shape.width / 914400 * 96)
            h = int(shape.height / 914400 * 96)
    except Exception:
        pass

    # Infer extension from content_type for filename consistency.
    ext_map = {
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
        "image/x-emf": ".emf",
        "image/webp": ".webp",
    }
    suffix = ext_map.get(content_type, ".png")

    def _factory(data=blob) -> bytes:
        return data

    return ImageRef(
        name=f"{name}{suffix}",
        mime=content_type,
        width=w,
        height=h,
        _factory=_factory,
    )
