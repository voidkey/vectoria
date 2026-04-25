"""Native .pptx parser via python-pptx.

Produces both text and images from a single slide walk. Speaker-
notes text is emitted inline (docling used to drop these; that was
the original reason to build W4-b's image-extractor override seam
around it — collapsed back into the parser in W6-6 since the override
had exactly one implementation and no other consumers).

Output structure
----------------
Each slide becomes an ``## Slide N: Title`` section. Body text
(title + text frames + cell text from simple tables) follows. If the
slide has a notes slide, a nested ``### Notes`` block carries the
notes text — same-paragraph preservation, not merged in-line.

Images
------
Both slide body shapes and notes-slide shapes are walked. Picture
shapes (MSO_SHAPE_TYPE.PICTURE) emit ``ImageRef`` entries with
``BytesFactory`` (picklable — parser runs under isolation). Names:
  slide_{idx:03d}_img_{n}          — body picture
  slide_{idx:03d}_notes_img_{n}    — speaker-notes picture
so downstream DB rows + S3 keys are greppable to "which slide did
this come from, and was it in notes".
"""
from __future__ import annotations

import asyncio
import io
import logging
from pathlib import Path

from config import get_settings
from parsers.base import BaseParser, ParseResult
from parsers.convert import LEGACY_FORMAT_MAP, convert_legacy_format
from parsers.image_ref import BytesFactory, ImageRef
from parsers.isolation import run_isolated

logger = logging.getLogger(__name__)

_MIME_EXT_MAP = {
    "image/png":   ".png",
    "image/jpeg":  ".jpg",
    "image/gif":   ".gif",
    "image/bmp":   ".bmp",
    "image/tiff":  ".tiff",
    "image/x-emf": ".emf",
    "image/webp":  ".webp",
}


class PptxParser(BaseParser):
    engine_name = "pptx-native"
    supported_types = [".pptx", ".ppt"]

    @classmethod
    def is_available(cls) -> bool:
        try:
            import pptx  # noqa: F401
        except ImportError:
            return False
        return True

    async def parse(
        self, source: bytes | str, filename: str = "", **kwargs,
    ) -> ParseResult:
        cfg = get_settings()
        if not cfg.parser_isolation:
            return await asyncio.get_running_loop().run_in_executor(
                None, self._parse_sync, source, filename,
            )
        return await run_isolated(
            _pptx_parse_worker, source, filename,
            timeout=cfg.parser_timeout, tier="fast",
        )

    def _parse_sync(self, source: bytes | str, filename: str) -> ParseResult:
        from pptx import Presentation

        raw = source if isinstance(source, bytes) else source.encode()

        suffix = Path(filename).suffix.lower()
        if suffix in LEGACY_FORMAT_MAP:
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                tmp.write(raw)
                tmp_path = tmp.name
            try:
                converted_path = convert_legacy_format(tmp_path, suffix)
                try:
                    with open(converted_path, "rb") as f:
                        raw = f.read()
                finally:
                    Path(converted_path).unlink(missing_ok=True)
            finally:
                Path(tmp_path).unlink(missing_ok=True)

        try:
            prs = Presentation(io.BytesIO(raw))
        except Exception:
            logger.exception("python-pptx parse failed for %s", filename)
            return ParseResult(content="", title=Path(filename).stem)

        lines: list[str] = []
        image_refs: list[ImageRef] = []
        deck_title = _extract_deck_title(prs) or Path(filename).stem
        lines.append(f"# {deck_title}")
        lines.append("")

        for idx, slide in enumerate(prs.slides, start=1):
            slide_idx_0 = idx - 1
            slide_title = _extract_slide_title(slide)
            header = f"## Slide {idx}"
            if slide_title:
                header = f"{header}: {slide_title}"
            lines.append(header)
            lines.append("")

            body = _extract_slide_body(slide, skip_title=slide_title)
            if body:
                lines.append(body)
                lines.append("")

            # Body images.
            _collect_picture_refs(
                slide.shapes, image_refs,
                name_prefix=f"slide_{slide_idx_0:03d}_img",
            )

            # Speaker notes — the docling-drop gap the native path fills.
            if slide.has_notes_slide:
                notes_text = _extract_notes_text(slide.notes_slide)
                if notes_text:
                    lines.append("### Notes")
                    lines.append("")
                    lines.append(notes_text)
                    lines.append("")
                _collect_picture_refs(
                    slide.notes_slide.shapes, image_refs,
                    name_prefix=f"slide_{slide_idx_0:03d}_notes_img",
                )

        content = "\n".join(lines).strip() + "\n"

        return ParseResult(
            content=content,
            title=deck_title,
            image_refs=image_refs,
        )


def _extract_deck_title(prs) -> str:
    """Deck title from the core-properties metadata (Title field)."""
    try:
        title = (prs.core_properties.title or "").strip()
    except Exception:
        return ""
    return title


def _extract_slide_title(slide) -> str:
    """Title placeholder text, or empty string."""
    try:
        if slide.shapes.title is not None:
            return (slide.shapes.title.text_frame.text or "").strip()
    except Exception:
        pass
    return ""


def _extract_slide_body(slide, *, skip_title: str = "") -> str:
    """Concatenate non-title text frames and simple table cell text.

    Preserves paragraph boundaries (one per line). Skips the title
    frame (already emitted as the ``## Slide N: Title`` header) and
    empty frames. Groups recurse.
    """
    parts: list[str] = []

    def _walk_shapes(shapes):
        for shape in shapes:
            if shape.shape_type == 6:  # GROUP
                _walk_shapes(shape.shapes)
                continue
            # Title placeholder already in header — skip.
            if shape == getattr(slide.shapes, "title", None):
                continue
            if shape.has_text_frame:
                txt = _text_frame_to_str(shape.text_frame)
                if txt and txt != skip_title:
                    parts.append(txt)
            if shape.has_table:
                parts.append(_table_to_markdown(shape.table))

    _walk_shapes(slide.shapes)
    return "\n\n".join(p for p in parts if p)


def _text_frame_to_str(tf) -> str:
    """One paragraph per line, blank paragraphs dropped."""
    paras = []
    for p in tf.paragraphs:
        txt = "".join(r.text for r in p.runs).strip()
        if txt:
            paras.append(txt)
    return "\n".join(paras)


def _table_to_markdown(table) -> str:
    """Pipe-style markdown table. Simple: no cell merging awareness —
    merged cells emit the same text per physical row/column (which is
    what python-pptx gives us). RAG use doesn't care about merge
    accuracy; we care about the text content.
    """
    rows = []
    for row in table.rows:
        cells = [
            (cell.text_frame.text or "").strip().replace("\n", " ").replace("|", "\\|")
            for cell in row.cells
        ]
        rows.append("| " + " | ".join(cells) + " |")
    if len(rows) >= 1:
        header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
        rows.insert(1, header_sep)
    return "\n".join(rows)


def _extract_notes_text(notes_slide) -> str:
    """Text from the notes_text_frame — skip the slide-number
    placeholder python-pptx sometimes leaves at the top of a notes
    frame when the deck hasn't been edited.
    """
    try:
        tf = notes_slide.notes_text_frame
    except Exception:
        return ""
    if tf is None:
        return ""
    return _text_frame_to_str(tf)


def _collect_picture_refs(
    shapes, out: list[ImageRef], *, name_prefix: str,
) -> None:
    """Walk a shape tree, appending ImageRef entries for every Picture.

    Recurses into GroupShapes. Uses ``BytesFactory`` so the ImageRef
    round-trips through the process-pool boundary — parsers run under
    ``parser_isolation``, and nested closures capturing pptx objects
    would fail to unpickle in the parent.
    """
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            _collect_picture_refs(shape.shapes, out, name_prefix=name_prefix)
            continue
        # Picture.image is a property that raises ValueError("no embedded
        # image") when the picture placeholder has no rId yet (e.g. an
        # unfilled "Picture with Caption" layout). getattr-with-default
        # only swallows AttributeError, not ValueError, so that path used
        # to bubble out and kill the whole slide walk. Wrap the access
        # itself; AttributeError covers shapes without .image at all
        # (text frames, tables, lines, ...).
        try:
            image = shape.image
            blob = image.blob
            content_type = image.content_type or "image/png"
        except (AttributeError, ValueError):
            continue

        n = len(out)
        suffix = _MIME_EXT_MAP.get(content_type, ".png")
        # Width/height from the slide geometry when available (EMU:
        # 914400 per inch; px assumes 96 DPI).
        w = h = None
        try:
            if shape.width and shape.height:
                w = int(shape.width / 914400 * 96)
                h = int(shape.height / 914400 * 96)
        except Exception:
            pass

        out.append(ImageRef(
            name=f"{name_prefix}_{n}{suffix}",
            mime=content_type,
            width=w,
            height=h,
            _factory=BytesFactory(blob),
        ))


def _pptx_parse_worker(source: bytes | str, filename: str) -> ParseResult:
    """Subprocess entry for PptxParser."""
    return PptxParser()._parse_sync(source, filename)
