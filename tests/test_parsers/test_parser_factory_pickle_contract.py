"""Contract test: every parser that builds ``ImageRef._factory`` must
produce refs that survive pickle.

Why
---
Parsers run under ``parser_isolation`` in production, which wraps them
in a ``ProcessPoolExecutor``. The returned ``ParseResult`` crosses the
process boundary via pickle. Nested ``def _factory`` closures serialize
by qualname (``ParserClass._method.<locals>._factory``) and fail to
unpickle in the parent with ``Can't get local object ...``, which
collapses the entire upload task into the DLQ.

This file pins down the picklability contract for each of the four
parser paths that build factories, so any future parser that reintroduces
a nested closure pattern fails here instead of silently in prod.
"""
from __future__ import annotations

import base64
import io
import pickle
from unittest.mock import MagicMock, patch

import pytest

from parsers.image_ref import ImageRef


def _assert_refs_pickle(refs: list[ImageRef]) -> None:
    """Pickle each ref and confirm materialize() still returns bytes.

    Going through ``pickle.dumps`` → ``pickle.loads`` reproduces the
    exact path that ``ProcessPoolExecutor`` takes when returning a
    ``ParseResult`` from a worker to the parent.
    """
    assert refs, "fixture must produce at least one ref or the test is vacuous"
    for ref in refs:
        rt = pickle.loads(pickle.dumps(ref))
        # materialize must still work after round-trip
        data = rt.materialize()
        assert isinstance(data, bytes) and len(data) > 0


def test_docling_extract_image_refs_produces_picklable_refs():
    """Docling path: ``_extract_image_refs`` must eagerly materialize
    PNG bytes because the closure-captured ``picture`` / ``document``
    objects are not picklable.
    """
    from parsers.docling_parser import DoclingParser

    # Minimal PIL image so size + PNG encode both work.
    try:
        from PIL import Image
    except ImportError:
        pytest.skip("PIL not available")
    img = Image.new("RGB", (4, 4), color="red")

    pic = MagicMock()
    pic.get_image.return_value = img
    doc = MagicMock()
    doc.pictures = [pic]
    result = MagicMock()
    result.document = doc

    refs = DoclingParser()._extract_image_refs(result)
    _assert_refs_pickle(refs)


def test_mineru_build_image_refs_produces_picklable_refs():
    """MinerU path: base64 payload factory must hold only the payload
    string, not a nested closure or an ``re.Match`` object."""
    from parsers.mineru_parser import MinerUParser

    raw = b"tiny-png-placeholder-bytes"
    b64 = base64.b64encode(raw).decode()

    refs = MinerUParser()._build_image_refs({
        "plain.png": b64,
        "datauri.png": f"data:image/png;base64,{b64}",
    })
    _assert_refs_pickle(refs)
    # Sanity: both paths decode back to the original bytes.
    for ref in refs:
        assert pickle.loads(pickle.dumps(ref)).materialize() == raw


def test_docx_parser_produces_picklable_refs_for_docx_with_image():
    """DocxParser path: mammoth's ``_image_handler`` must build refs
    whose factory is module-level (``BytesFactory``), not a nested def."""
    try:
        from docx import Document
        from docx.shared import Inches  # noqa: F401 — availability check
    except ImportError:
        pytest.skip("python-docx not available")
    from parsers.docx_parser import DocxParser

    # Build a minimal .docx with a single 1×1 PNG embedded.
    try:
        from PIL import Image
    except ImportError:
        pytest.skip("PIL not available")

    png_buf = io.BytesIO()
    Image.new("RGB", (2, 2), color="blue").save(png_buf, format="PNG")
    png_buf.seek(0)

    doc = Document()
    doc.add_paragraph("hello")
    doc.add_picture(png_buf)
    out = io.BytesIO()
    doc.save(out)

    result = DocxParser()._parse_sync(out.getvalue(), "test.docx")
    _assert_refs_pickle(result.image_refs)


@pytest.mark.asyncio
async def test_pptx_image_extractor_produces_picklable_refs():
    """PptxImageExtractor path: ``_make_ref`` in ``parsers._pptx_images``
    must use ``BytesFactory`` rather than a nested def."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        pytest.skip("python-pptx not available")
    try:
        from PIL import Image
    except ImportError:
        pytest.skip("PIL not available")

    from parsers._pptx_images import PptxImageExtractor

    # Build a minimal .pptx with one slide holding one picture shape.
    png = io.BytesIO()
    Image.new("RGB", (8, 8), color="green").save(png, format="PNG")
    png.seek(0)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(png, Inches(1), Inches(1), Inches(1), Inches(1))
    out = io.BytesIO()
    prs.save(out)

    refs = await PptxImageExtractor().extract(out.getvalue(), filename="deck.pptx")
    _assert_refs_pickle(refs)


def test_parser_modules_have_no_nested_factory_defs():
    """Static safeguard: grep each parser module's source for the exact
    anti-pattern (``def _factory`` inside a method body). The fix
    replaced every such site with a module-level ``BytesFactory`` /
    ``Base64Factory``. A future parser reintroducing the nested form
    will fail here immediately rather than in prod's DLQ.

    Deliberately narrow: only ``def _factory(`` — we're not trying to
    ban every nested def in the package, just the specific
    unpicklable-factory footgun we just paid for.
    """
    import re
    from pathlib import Path

    parsers_dir = Path(__file__).resolve().parents[2] / "parsers"
    offenders = []
    for py in parsers_dir.rglob("*.py"):
        text = py.read_text(encoding="utf-8")
        # Match ``def _factory(`` indented (inside a function/method).
        # Module-level ``def _factory(`` would be picklable and is allowed,
        # but we don't currently have any, and nested ones are what bite us.
        for m in re.finditer(r"^(?P<indent>[ \t]+)def _factory\(", text,
                             flags=re.MULTILINE):
            if m.group("indent"):
                offenders.append(str(py.relative_to(parsers_dir.parent)))
                break
    assert not offenders, (
        "Nested 'def _factory(...)' found in parser modules "
        f"{offenders}. These closures cannot be pickled across a "
        "ProcessPoolExecutor boundary (parser_isolation). Use "
        "parsers.image_ref.BytesFactory / Base64Factory instead."
    )
