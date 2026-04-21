"""PptxImageExtractor: slides + speaker notes coverage.

The differentiating feature vs docling is pulling images out of
speaker-notes slides. We build a real in-memory .pptx with
python-pptx, attach one image to the main slide and another to the
notes, then assert the extractor finds both.
"""
import io

import pytest
from PIL import Image

from parsers._pptx_images import PptxImageExtractor
from parsers.image_ref import ImageRef


def _png_bytes(color=(200, 50, 100)) -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", (64, 64), color).save(buf, format="PNG")
    return buf.getvalue()


def _make_pptx_with_body_images() -> bytes:
    """Build a pptx in memory: 2 slides, each with one body image.

    python-pptx's ``NotesSlideShapes`` doesn't expose ``add_picture``,
    so we can't construct a notes-slide image programmatically. The
    body-images path is the bulk of the logic (shape tree walk +
    ``.image.blob`` + EMU→px math) and is exercised here; the
    notes-slide path uses the exact same ``_collect_shape_images``
    helper and is covered separately via a mock-based test that
    asserts ``notes_slide.shapes`` actually gets walked.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    for color in [(200, 50, 100), (50, 50, 200)]:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        img = io.BytesIO(_png_bytes(color=color))
        slide.shapes.add_picture(
            img, Inches(1), Inches(1),
            width=Inches(2), height=Inches(2),
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# match()
# ---------------------------------------------------------------------------

def test_match_claims_pptx_by_extension():
    e = PptxImageExtractor()
    assert e.match(ext=".pptx")
    assert not e.match(ext=".docx")
    assert not e.match(ext=".pdf")


def test_match_claims_pptx_by_mime():
    e = PptxImageExtractor()
    assert e.match(
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    assert not e.match(mime="image/png")


# ---------------------------------------------------------------------------
# extract()
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_extract_returns_body_images_from_slides():
    """Body-image path against a real pptx: 2 slides × 1 picture each."""
    pptx_bytes = _make_pptx_with_body_images()
    refs = await PptxImageExtractor().extract(pptx_bytes, filename="deck.pptx")

    assert len(refs) == 2

    # Names encode slide index and note-vs-body distinction, greppable
    # downstream in the document_images table + S3 keys.
    names = [r.name for r in refs]
    assert any("slide_000_img_0" in n for n in names)
    assert any("slide_001_img_0" in n for n in names)
    # No "notes" in these names — there were no notes images to find.
    assert not any("notes" in n for n in names)


@pytest.mark.asyncio
async def test_extract_walks_notes_slide_shapes():
    """Signature-capability guard: the extractor must walk
    ``slide.notes_slide.shapes`` when ``has_notes_slide`` is true. We
    can't programmatically add a picture to a notes slide via
    python-pptx (``NotesSlideShapes`` lacks ``add_picture``), so we
    mock the Presentation to feed the extractor one notes-slide
    picture and assert it surfaces as ``..._notes_img_0``.
    """
    from unittest.mock import MagicMock, patch

    # Shape on main slide
    body_shape = MagicMock(shape_type=13)  # MSO_SHAPE_TYPE.PICTURE
    body_shape.image = MagicMock(blob=b"BODY", content_type="image/png")
    body_shape.width = 914400 * 2
    body_shape.height = 914400 * 2

    # Shape on notes slide
    notes_shape = MagicMock(shape_type=13)
    notes_shape.image = MagicMock(blob=b"NOTES", content_type="image/png")
    notes_shape.width = 914400
    notes_shape.height = 914400

    slide = MagicMock(has_notes_slide=True)
    slide.shapes = [body_shape]
    slide.notes_slide.shapes = [notes_shape]

    fake_prs = MagicMock()
    fake_prs.slides = [slide]

    with patch("pptx.Presentation", return_value=fake_prs):
        refs = await PptxImageExtractor().extract(
            b"fake pptx bytes", filename="deck.pptx",
        )

    names = [r.name for r in refs]
    assert names == ["slide_000_img_0.png", "slide_000_notes_img_0.png"]
    # And the bytes carry the distinction — a regression where body
    # and notes bytes got swapped would leak here.
    body_ref = next(r for r in refs if "notes" not in r.name)
    notes_ref = next(r for r in refs if "notes" in r.name)
    assert body_ref.materialize() == b"BODY"
    assert notes_ref.materialize() == b"NOTES"


@pytest.mark.asyncio
async def test_extract_materialised_bytes_round_trip_through_pil():
    """Each ref's bytes_factory produces valid image bytes — validates
    we aren't mangling blobs on the way out.
    """
    pptx_bytes = _make_pptx_with_body_images()
    refs = await PptxImageExtractor().extract(pptx_bytes, filename="deck.pptx")

    for r in refs:
        data = r.materialize()
        assert len(data) > 0
        with Image.open(io.BytesIO(data)) as img:
            assert img.size == (64, 64)


@pytest.mark.asyncio
async def test_extract_returns_empty_for_malformed_bytes():
    """Invalid pptx bytes must log + return [] rather than raise."""
    refs = await PptxImageExtractor().extract(b"not a zip", filename="x.pptx")
    assert refs == []


@pytest.mark.asyncio
async def test_extract_handles_pptx_with_no_images():
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    buf = io.BytesIO()
    prs.save(buf)

    refs = await PptxImageExtractor().extract(buf.getvalue(), filename="x.pptx")
    assert refs == []


# ---------------------------------------------------------------------------
# Registration — the plugin is actually picked up at import time
# ---------------------------------------------------------------------------

def test_plugin_registered_at_parsers_import():
    from parsers import image_extractor as ie
    # After ``import parsers`` the registry must have a PptxImageExtractor
    # available for .pptx lookups, otherwise the integration with
    # worker/handlers.py won't fire.
    ext = ie.find_image_extractor(ext=".pptx")
    assert ext is not None
    assert isinstance(ext, PptxImageExtractor)
