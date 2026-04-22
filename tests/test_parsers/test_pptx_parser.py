"""PptxParser (python-pptx) — native .pptx text + image extraction.

Speaker notes (text and images) are the docling-drop gap the native
parser fills. Guards:
  * each slide emits as ``## Slide N: Title`` section
  * notes-slide text is nested as ``### Notes``
  * tables render as pipe-markdown
  * empty frames / blank paragraphs are dropped
  * body pictures surface in ``image_refs`` with
    ``slide_NNN_img_M`` naming
  * notes-slide pictures surface as ``slide_NNN_notes_img_M``
  * registry picks pptx-native over everything else
"""
import io

import pytest


def _build_pptx_with_notes() -> bytes:
    """Build a 2-slide deck with title, body, and speaker notes."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.core_properties.title = "Test Deck"

    # Slide 0: title + body + notes
    s0 = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    s0.shapes.title.text = "Opening"
    s0.placeholders[1].text = "Body paragraph one."
    s0.notes_slide.notes_text_frame.text = (
        "Presenter says: start with a question."
    )

    # Slide 1: title + body, no notes
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Closing"
    s1.placeholders[1].text = "Bullet body"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture(autouse=True)
def _disable_isolation(monkeypatch):
    from config import get_settings
    monkeypatch.setattr(get_settings(), "parser_isolation", False)


# ---------------------------------------------------------------------------
# Engine metadata
# ---------------------------------------------------------------------------

def test_engine_name_and_supported_types():
    from parsers.pptx_parser import PptxParser
    assert PptxParser.engine_name == "pptx-native"
    assert ".pptx" in PptxParser.supported_types
    assert ".ppt" in PptxParser.supported_types


def test_is_available_with_dep_present():
    from parsers.pptx_parser import PptxParser
    assert PptxParser.is_available()


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_parse_emits_slide_sections_with_titles():
    from parsers.pptx_parser import PptxParser
    refs_bytes = _build_pptx_with_notes()
    result = await PptxParser().parse(refs_bytes, filename="deck.pptx")

    assert "# Test Deck" in result.content
    assert "## Slide 1: Opening" in result.content
    assert "## Slide 2: Closing" in result.content


@pytest.mark.asyncio
async def test_parse_includes_speaker_notes_under_notes_heading():
    """The signature capability over docling: notes text surfaces."""
    from parsers.pptx_parser import PptxParser
    result = await PptxParser().parse(
        _build_pptx_with_notes(), filename="deck.pptx",
    )

    assert "### Notes" in result.content
    assert "start with a question" in result.content
    # Only slide 0 had notes — heading should appear exactly once.
    assert result.content.count("### Notes") == 1


@pytest.mark.asyncio
async def test_parse_includes_body_paragraphs():
    from parsers.pptx_parser import PptxParser
    result = await PptxParser().parse(
        _build_pptx_with_notes(), filename="deck.pptx",
    )
    assert "Body paragraph one" in result.content
    assert "Bullet body" in result.content


@pytest.mark.asyncio
async def test_parse_title_from_core_properties():
    from parsers.pptx_parser import PptxParser
    result = await PptxParser().parse(
        _build_pptx_with_notes(), filename="x.pptx",
    )
    assert result.title == "Test Deck"


@pytest.mark.asyncio
async def test_parse_title_falls_back_to_filename_when_no_core_title():
    from parsers.pptx_parser import PptxParser
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    buf = io.BytesIO()
    prs.save(buf)

    result = await PptxParser().parse(buf.getvalue(), filename="untitled.pptx")
    assert result.title == "untitled"


# ---------------------------------------------------------------------------
# Image path: body + notes pictures in a single slide walk (W6-6
# merged the old PptxImageExtractor back into the parser)
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_parse_produces_body_picture_image_refs():
    """Body Picture shapes must surface in image_refs with the
    ``slide_NNN_img_M`` naming convention.
    """
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image
    from parsers.pptx_parser import PptxParser

    png = io.BytesIO()
    Image.new("RGB", (8, 8), color="green").save(png, format="PNG")
    png.seek(0)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_picture(png, Inches(1), Inches(1), Inches(1), Inches(1))
    out = io.BytesIO()
    prs.save(out)

    result = await PptxParser().parse(out.getvalue(), filename="deck.pptx")
    assert len(result.image_refs) == 1
    assert result.image_refs[0].name.startswith("slide_000_img_0")
    # Dimensions propagated from the slide geometry (1 inch × 96 DPI).
    assert result.image_refs[0].width == 96
    assert result.image_refs[0].height == 96


@pytest.mark.asyncio
async def test_parse_returns_empty_image_refs_on_deck_without_pictures():
    """The text-only fixture has no pictures, so image_refs is empty.
    Guards against accidental ghost refs.
    """
    from parsers.pptx_parser import PptxParser
    result = await PptxParser().parse(
        _build_pptx_with_notes(), filename="deck.pptx",
    )
    assert result.image_refs == []


# ---------------------------------------------------------------------------
# Error handling
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_parse_handles_malformed_bytes():
    from parsers.pptx_parser import PptxParser
    result = await PptxParser().parse(b"not a pptx", filename="bad.pptx")
    assert result.content == ""
    assert result.image_refs == []


# ---------------------------------------------------------------------------
# Registry dispatch
# ---------------------------------------------------------------------------

def test_registry_picks_pptx_native_over_docling():
    from parsers.registry import registry
    engine = registry.auto_select(filename="deck.pptx")
    assert engine == "pptx-native", (
        f"expected pptx-native, got {engine!r}; "
        "_EXT_PREFERENCE ordering probably regressed"
    )
