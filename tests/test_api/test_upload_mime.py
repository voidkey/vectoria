"""Magic-byte verification on upload — mime_sniff module + endpoint gate."""
import pytest
from datetime import datetime
from unittest.mock import patch, AsyncMock, MagicMock

from api.mime_sniff import detect_family, check_mime


# Smallest-known magic-byte headers for the formats we care about.
_PDF_HEAD = b"%PDF-1.4\n"
_PNG_HEAD = b"\x89PNG\r\n\x1a\n" + b"\x00" * 16
_JPG_HEAD = b"\xff\xd8\xff\xe0\x00\x10JFIF\x00\x01"
_ZIP_HEAD = b"PK\x03\x04" + b"\x00" * 26
_EXE_HEAD = b"MZ\x90\x00\x03\x00\x00\x00"


def test_detect_family_pdf():
    assert detect_family(_PDF_HEAD) == "pdf"


def test_detect_family_png():
    assert detect_family(_PNG_HEAD) == "image"


def test_detect_family_jpg():
    assert detect_family(_JPG_HEAD) == "image"


def test_detect_family_office_zip():
    # Plain zip head may be ambiguous (no central dir yet).
    result = detect_family(_ZIP_HEAD)
    assert result in (None, "office-doc", "office-sheet", "office-slide")


def test_check_mime_ok_on_match():
    ok, detected = check_mime("ok.pdf", _PDF_HEAD)
    assert ok is True
    assert detected == "pdf"


def test_check_mime_rejects_exe_as_pdf():
    ok, detected = check_mime("evil.pdf", _EXE_HEAD)
    assert ok is False


def test_check_mime_passes_when_undetectable():
    # Plain text — no magic signature — should pass through ambiguously.
    ok, detected = check_mime("note.md", b"# hello\n\nbody")
    assert ok is True
    assert detected in (None, "text")


def test_check_mime_rejects_pdf_as_docx():
    ok, detected = check_mime("oops.docx", _PDF_HEAD)
    assert ok is False
    assert detected == "pdf"


def test_check_mime_rejects_exe_as_exe():
    """An executable uploaded as .exe must be rejected even though .exe
    isn't in EXT_FAMILIES — blocked family wins over pass-through."""
    ok, detected = check_mime("malware.exe", _EXE_HEAD)
    assert ok is False
    assert detected == "executable"


def test_check_mime_accepts_real_pptx_against_multimatch_ambiguity():
    """Regression: puremagic returns ~8 equal-confidence guesses for any
    OOXML zip (docx/pptx/xlsx variants all share the same magic). The
    old detect_family() returned the *first* match's family, which for
    a real .pptx happened to be office-doc — so legitimate .pptx
    uploads were rejected with MIME_MISMATCH. check_mime now allows as
    long as the claim shows up *anywhere* in puremagic's candidate set.
    """
    import io
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[1]).shapes.title.text = "hi"
    buf = io.BytesIO(); prs.save(buf)
    head = buf.getvalue()[:2048]

    ok_pptx, fam_pptx = check_mime("deck.pptx", head)
    assert ok_pptx is True, f"pptx rejected, detected={fam_pptx}"

    # Symmetry: the same OOXML head, claimed as .docx, must also pass —
    # puremagic can't tell docx from pptx from this head.
    ok_docx, _ = check_mime("paper.docx", head)
    assert ok_docx is True


def test_check_mime_still_rejects_pdf_claimed_as_pptx():
    """Sanity: the multi-match relaxation must NOT allow real
    cross-family forgeries. A PDF head claimed as .pptx still fails.
    """
    ok, detected = check_mime("forgery.pptx", _PDF_HEAD)
    assert ok is False
    assert detected == "pdf"


def test_detect_families_returns_all_ooxml_candidates():
    from api.mime_sniff import detect_families
    import io
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[1])
    buf = io.BytesIO(); prs.save(buf)
    fams = detect_families(buf.getvalue()[:2048])
    # Real pptx surfaces as multiple OOXML sub-families simultaneously.
    assert "office-slide" in fams
    assert len(set(fams)) >= 2  # at minimum slide + one of doc/sheet


def test_check_mime_rejects_exe_as_unknown_ext():
    """Same hole, confirmed for any unknown/odd extension."""
    ok, detected = check_mime("weird.xyz", _EXE_HEAD)
    assert ok is False
    assert detected == "executable"


# -- Endpoint-level tests --------------------------------------------------
# Fixtures (client, kb) reuse the pattern from
# tests/test_api/test_upload_limits_and_dedup.py — the `client` fixture is
# defined in tests/conftest.py and shared automatically.


def _make_session_miss():
    """Return a mock async session that returns None on execute (dedup miss)."""
    session = AsyncMock()
    miss_result = MagicMock()
    miss_result.scalar_one_or_none.return_value = None
    session.execute = AsyncMock(return_value=miss_result)

    def _refresh(obj):
        obj.created_at = datetime(2026, 4, 15, 21, 0, 0)

    session.add = MagicMock()
    session.commit = AsyncMock()
    session.refresh = AsyncMock(side_effect=_refresh)
    return session


@pytest.mark.asyncio
async def test_upload_pdf_bytes_as_pdf_accepted(client):
    """Valid PDF magic bytes with .pdf extension → 201."""
    pdf_content = _PDF_HEAD + b"%more pdf content here"
    session = _make_session_miss()

    with (
        patch("api.routes.documents._validate_kb", new=AsyncMock()),
        patch("api.routes.documents.get_storage") as mock_storage,
        patch("api.routes.documents.registry") as mock_registry,
        patch("api.routes.documents.get_session") as mock_sess,
    ):
        mock_storage.return_value = AsyncMock()
        mock_registry.auto_select.return_value = "markitdown"
        mock_sess.return_value.__aenter__.return_value = session

        resp = await client.post(
            "/v1/knowledgebases/kb-x/documents/file",
            files={"file": ("ok.pdf", pdf_content, "application/pdf")},
        )

    assert resp.status_code == 201, resp.text


@pytest.mark.asyncio
async def test_upload_exe_as_pdf_rejected(client):
    """PE executable bytes with .pdf extension → 400 mime_mismatch."""
    exe_content = _EXE_HEAD + b"\x00" * 64

    with (
        patch("api.routes.documents._validate_kb", new=AsyncMock()),
        patch("api.routes.documents.get_storage") as mock_storage,
        patch("api.routes.documents.registry") as mock_registry,
    ):
        mock_storage.return_value = AsyncMock()
        mock_registry.auto_select.return_value = "markitdown"

        resp = await client.post(
            "/v1/knowledgebases/kb-x/documents/file",
            files={"file": ("evil.pdf", exe_content, "application/pdf")},
        )

    assert resp.status_code == 400, resp.text
    body = resp.json()
    assert body["code"] == 1207  # MIME_MISMATCH


@pytest.mark.asyncio
async def test_upload_pdf_as_docx_rejected(client):
    """PDF magic bytes with .docx extension → 400 (cross-family mismatch)."""
    pdf_content = _PDF_HEAD + b"%more content"

    with (
        patch("api.routes.documents._validate_kb", new=AsyncMock()),
        patch("api.routes.documents.get_storage") as mock_storage,
        patch("api.routes.documents.registry") as mock_registry,
    ):
        mock_storage.return_value = AsyncMock()
        mock_registry.auto_select.return_value = "markitdown"

        resp = await client.post(
            "/v1/knowledgebases/kb-x/documents/file",
            files={"file": ("oops.docx", pdf_content, "application/pdf")},
        )

    assert resp.status_code == 400, resp.text
    body = resp.json()
    assert body["code"] == 1207  # MIME_MISMATCH


@pytest.mark.asyncio
async def test_upload_markdown_accepted_ambiguous(client):
    """Plain text with .md extension → 201 (ambiguous detection passes through)."""
    md_content = b"# hello\n\nThis is a markdown document.\n"
    session = _make_session_miss()

    with (
        patch("api.routes.documents._validate_kb", new=AsyncMock()),
        patch("api.routes.documents.get_storage") as mock_storage,
        patch("api.routes.documents.registry") as mock_registry,
        patch("api.routes.documents.get_session") as mock_sess,
    ):
        mock_storage.return_value = AsyncMock()
        mock_registry.auto_select.return_value = "markitdown"
        mock_sess.return_value.__aenter__.return_value = session

        resp = await client.post(
            "/v1/knowledgebases/kb-x/documents/file",
            files={"file": ("note.md", md_content, "text/markdown")},
        )

    assert resp.status_code == 201, resp.text


@pytest.mark.asyncio
async def test_upload_non_strict_allows_mismatch(monkeypatch, client):
    """With strict_mime_check=False, exe-as-pdf → 201 (but counter still fires)."""
    from config import get_settings
    monkeypatch.setattr(get_settings(), "strict_mime_check", False)

    exe_content = _EXE_HEAD + b"\x00" * 64
    session = _make_session_miss()

    with (
        patch("api.routes.documents._validate_kb", new=AsyncMock()),
        patch("api.routes.documents.get_storage") as mock_storage,
        patch("api.routes.documents.registry") as mock_registry,
        patch("api.routes.documents.get_session") as mock_sess,
        patch("infra.metrics.UPLOAD_MIME_MISMATCH_TOTAL") as mock_counter,
    ):
        mock_storage.return_value = AsyncMock()
        mock_registry.auto_select.return_value = "markitdown"
        mock_sess.return_value.__aenter__.return_value = session

        mock_labels = MagicMock()
        mock_counter.labels.return_value = mock_labels

        resp = await client.post(
            "/v1/knowledgebases/kb-x/documents/file",
            files={"file": ("evil.pdf", exe_content, "application/pdf")},
        )

    assert resp.status_code == 201, resp.text
    # Counter must have been incremented even in non-strict mode.
    mock_counter.labels.assert_called_once()
    mock_labels.inc.assert_called_once()


@pytest.mark.asyncio
async def test_upload_exe_as_exe_rejected(client):
    """Executable with true .exe name still rejected — blocked family wins."""
    exe_content = _EXE_HEAD + b"\x00" * 64

    with (
        patch("api.routes.documents._validate_kb", new=AsyncMock()),
        patch("api.routes.documents.get_storage") as mock_storage,
        patch("api.routes.documents.registry") as mock_registry,
    ):
        mock_storage.return_value = AsyncMock()
        mock_registry.auto_select.return_value = "markitdown"

        resp = await client.post(
            "/v1/knowledgebases/kb-x/documents/file",
            files={"file": ("actually.exe", exe_content, "application/octet-stream")},
        )

    assert resp.status_code == 400, resp.text
    body = resp.json()
    assert body["code"] == 1207  # MIME_MISMATCH


# ---------------------------------------------------------------------------
# upload_rejected log + counter — operator visibility for 4xx rejects
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_mime_reject_emits_upload_rejected_log_and_counter(client, caplog):
    """Strict-mode MIME reject must (a) WARN-log ``upload_rejected ...``
    so operators can grep for it, and (b) bump
    ``UPLOAD_REJECTED_TOTAL{reason="mime_mismatch", claimed_ext=".pdf"}``
    so a single alert rule can fire on cumulative rejection rate.
    The HTTP access log only carries the status code — without this
    pair you can't tell which file got rejected.
    """
    import logging
    exe_content = _EXE_HEAD + b"\x00" * 64

    with (
        patch("api.routes.documents._validate_kb", new=AsyncMock()),
        patch("api.routes.documents.get_storage") as mock_storage,
        patch("api.routes.documents.registry") as mock_registry,
        patch("infra.metrics.UPLOAD_REJECTED_TOTAL") as mock_rejected,
    ):
        mock_storage.return_value = AsyncMock()
        mock_registry.auto_select.return_value = "markitdown"
        mock_labels = MagicMock()
        mock_rejected.labels.return_value = mock_labels

        with caplog.at_level(logging.WARNING, logger="api.routes.documents"):
            resp = await client.post(
                "/v1/knowledgebases/kb-evil/documents/file",
                files={"file": ("evil.pdf", exe_content, "application/pdf")},
            )

    assert resp.status_code == 400
    # Counter: bounded labels only.
    mock_rejected.labels.assert_called_once_with(
        reason="mime_mismatch", claimed_ext=".pdf",
    )
    mock_labels.inc.assert_called_once()
    # Log: the unbounded specifics live here so we never blow up Prom
    # cardinality but still recover the filename for forensics.
    rejects = [r for r in caplog.records if "upload_rejected" in r.getMessage()]
    assert len(rejects) == 1
    msg = rejects[0].getMessage()
    assert "kb=kb-evil" in msg
    assert "filename=evil.pdf" in msg
    assert "reason=mime_mismatch" in msg
    assert "detected=executable" in msg


@pytest.mark.asyncio
async def test_too_large_reject_emits_upload_rejected_log_and_counter(client, caplog):
    """Same pair (log + counter) for the 413 size-gate rejection path,
    so all 4xx upload rejects are uniformly observable from one place.
    """
    import logging
    from config import get_settings
    limit = get_settings().max_upload_bytes
    oversized = b"a" * (limit + 1)

    with (
        patch("api.routes.documents._validate_kb", new=AsyncMock()),
        patch("api.routes.documents.get_storage") as mock_storage,
        patch("api.routes.documents.registry") as mock_reg,
        patch("infra.metrics.UPLOAD_REJECTED_TOTAL") as mock_rejected,
    ):
        mock_storage.return_value = AsyncMock()
        mock_reg.auto_select.return_value = "markitdown"
        mock_labels = MagicMock()
        mock_rejected.labels.return_value = mock_labels

        with caplog.at_level(logging.WARNING, logger="api.routes.documents"):
            resp = await client.post(
                "/v1/knowledgebases/kb-big/documents/file",
                files={"file": ("huge.bin", oversized, "application/octet-stream")},
            )

    assert resp.status_code == 413
    mock_rejected.labels.assert_called_once_with(
        reason="too_large", claimed_ext="other",  # .bin not in EXT_FAMILIES
    )
    mock_labels.inc.assert_called_once()
    rejects = [r for r in caplog.records if "upload_rejected" in r.getMessage()]
    assert len(rejects) == 1
    msg = rejects[0].getMessage()
    assert "kb=kb-big" in msg
    assert "filename=huge.bin" in msg
    assert "reason=too_large" in msg
    assert f"limit={limit}" in msg
